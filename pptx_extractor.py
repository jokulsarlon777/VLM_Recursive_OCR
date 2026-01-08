"""
PowerPoint OLE Object Extractor
Extracts embedded OLE objects (ppt/pptx) from PowerPoint files recursively

Supports multiple extraction methods:
1. ZIP-based extraction (fastest, works for standard PPTX)
2. COM automation extraction (works for corrupted/non-standard PPTX)
3. python-pptx OLE extraction (fallback)
"""
import os
import sys
import tempfile
import zipfile
import shutil
from pathlib import Path
from typing import List, Dict, Tuple, Optional
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.oxml import parse_xml
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Import COM only on Windows
if sys.platform == 'win32':
    try:
        import win32com.client
        import pythoncom
        COM_AVAILABLE = True
    except ImportError:
        COM_AVAILABLE = False
        logger.warning("win32com not available, COM extraction will be disabled")
else:
    COM_AVAILABLE = False


class PPTXExtractor:
    """Extract OLE objects from PowerPoint files"""

    def __init__(self, pptx_path: Path):
        """
        Initialize the extractor

        Args:
            pptx_path: Path to the PowerPoint file
        """
        self.pptx_path = Path(pptx_path)
        self.presentation = None

    def load_presentation(self) -> None:
        """Load the PowerPoint presentation"""
        try:
            self.presentation = Presentation(str(self.pptx_path))
            logger.info(f"Loaded presentation: {self.pptx_path.name}")
        except Exception as e:
            logger.error(f"Failed to load presentation: {e}")
            raise

    def extract_ole_objects(self) -> List[Dict]:
        """
        Extract all OLE objects from the presentation

        Returns:
            List of dictionaries containing OLE object information
        """
        if not self.presentation:
            self.load_presentation()

        ole_objects = []

        for slide_idx, slide in enumerate(self.presentation.slides, 1):
            logger.info(f"Processing slide {slide_idx}/{len(self.presentation.slides)}")

            for shape_idx, shape in enumerate(slide.shapes):
                # Check if shape contains OLE object
                if self._is_ole_object(shape):
                    ole_info = self._extract_ole_data(shape, slide_idx, shape_idx)
                    if ole_info:
                        ole_objects.append(ole_info)

        logger.info(f"Found {len(ole_objects)} OLE objects")
        return ole_objects

    def _is_ole_object(self, shape) -> bool:
        """
        Check if a shape is an OLE object

        Args:
            shape: PowerPoint shape object

        Returns:
            True if shape is an OLE object
        """
        try:
            # Check if shape has oleObject element
            if hasattr(shape, '_element'):
                xml = shape._element.xml
                return b'oleObject' in xml or b'embed' in xml
            return False
        except Exception as e:
            logger.debug(f"Error checking OLE object: {e}")
            return False

    def _extract_ole_data(self, shape, slide_idx: int, shape_idx: int) -> Dict:
        """
        Extract OLE object data from a shape

        Args:
            shape: PowerPoint shape containing OLE object
            slide_idx: Slide index
            shape_idx: Shape index within the slide

        Returns:
            Dictionary with OLE object information
        """
        try:
            # Get the relationship ID for the embedded object
            ole_elem = shape._element

            # Try to find the embedded object relationship
            for rel in shape.part.rels.values():
                if 'oleObject' in rel.reltype or 'package' in rel.reltype:
                    # Extract the blob data
                    blob_data = rel.target_part.blob

                    # Determine file extension
                    content_type = rel.target_part.content_type
                    ext = self._get_extension_from_content_type(content_type)

                    return {
                        'slide_idx': slide_idx,
                        'shape_idx': shape_idx,
                        'blob_data': blob_data,
                        'extension': ext,
                        'content_type': content_type,
                        'filename': f"ole_s{slide_idx}_sh{shape_idx}{ext}"
                    }

            return None

        except Exception as e:
            logger.error(f"Error extracting OLE data from slide {slide_idx}, shape {shape_idx}: {e}")
            return None

    def _get_extension_from_content_type(self, content_type: str) -> str:
        """
        Get file extension from content type

        Args:
            content_type: MIME type or content type string

        Returns:
            File extension with dot (e.g., '.pptx')
        """
        content_type_map = {
            'application/vnd.ms-powerpoint': '.ppt',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
            'application/vnd.ms-powerpoint.presentation.macroEnabled.12': '.pptm',
        }

        return content_type_map.get(content_type, '.bin')

    def save_ole_objects(self, ole_objects: List[Dict], output_dir: Path) -> List[Path]:
        """
        Save extracted OLE objects to files

        Args:
            ole_objects: List of OLE object dictionaries
            output_dir: Directory to save extracted files

        Returns:
            List of paths to saved files
        """
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        saved_files = []

        for ole_obj in ole_objects:
            try:
                # Only save PowerPoint files
                if ole_obj['extension'] in ['.ppt', '.pptx', '.pptm']:
                    output_path = output_dir / ole_obj['filename']

                    with open(output_path, 'wb') as f:
                        f.write(ole_obj['blob_data'])

                    saved_files.append(output_path)
                    logger.info(f"Saved OLE object to: {output_path}")

            except Exception as e:
                logger.error(f"Error saving OLE object {ole_obj['filename']}: {e}")

        return saved_files


def extract_embedded_pptx_from_zip(pptx_path: Path, output_dir: Path) -> List[Path]:
    """
    Extract embedded files directly from PPTX ZIP structure

    Args:
        pptx_path: Path to the source PowerPoint file
        output_dir: Directory to save extracted files

    Returns:
        List of paths to extracted PowerPoint files
    """
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    extracted_files = []

    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            # Look for embedded files in ppt/embeddings/ directory
            embedded_files = [f for f in zip_ref.namelist() if f.startswith('ppt/embeddings/')]

            logger.info(f"Found {len(embedded_files)} files in ppt/embeddings/")

            for idx, file_path in enumerate(embedded_files, 1):
                try:
                    # Extract the file
                    file_data = zip_ref.read(file_path)

                    # Determine file extension by checking file signature
                    ext = _detect_file_extension(file_data)

                    # Skip if not a PowerPoint file
                    if ext not in ['.ppt', '.pptx', '.pptm']:
                        logger.debug(f"Skipping non-PowerPoint file: {file_path} (detected: {ext})")
                        continue

                    # Generate output filename
                    output_filename = f"embedded_{idx}{ext}"
                    output_path = output_dir / output_filename

                    # Save the file
                    with open(output_path, 'wb') as f:
                        f.write(file_data)

                    # Validate the extracted file
                    if not _validate_extracted_file(output_path):
                        logger.warning(f"Extracted file validation failed: {output_filename}")
                        logger.warning(f"File may be corrupted or in OLE stream format - will be skipped during conversion")
                        # Keep the file but don't add to extracted_files list so it won't be processed
                        continue

                    extracted_files.append(output_path)
                    logger.info(f"Extracted embedded file: {output_filename}")

                except Exception as e:
                    logger.error(f"Error extracting {file_path}: {e}")

    except zipfile.BadZipFile:
        logger.error(f"Invalid PPTX file (not a valid ZIP): {pptx_path}")
    except Exception as e:
        logger.error(f"Error reading PPTX as ZIP: {e}")

    return extracted_files


def _validate_extracted_file(file_path: Path) -> bool:
    """
    Validate if an extracted PowerPoint file is valid

    Args:
        file_path: Path to the extracted file

    Returns:
        True if file appears valid, False otherwise
    """
    try:
        # Basic file size check
        file_size = file_path.stat().st_size
        if file_size < 512:  # Too small to be a valid PowerPoint file
            logger.debug(f"File too small: {file_size} bytes")
            return False

        # Read file header to validate
        with open(file_path, 'rb') as f:
            header = f.read(512)

        # Check for valid PPTX (ZIP) signature
        if header[:4] == b'PK\x03\x04':
            # Check if it has proper ZIP structure
            try:
                with zipfile.ZipFile(file_path, 'r') as zf:
                    namelist = zf.namelist()
                    # Valid PPTX should have [Content_Types].xml
                    if '[Content_Types].xml' in namelist:
                        return True
                    else:
                        logger.debug("Missing [Content_Types].xml in PPTX")
                        return False
            except zipfile.BadZipFile:
                logger.debug("Invalid ZIP structure")
                return False

        # Check for valid PPT (OLE2) signature
        elif header[:8] == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1':
            # This is an OLE2 file, should be valid
            # Additional validation could be done here
            return True

        else:
            logger.debug(f"Unknown file signature: {header[:8].hex()}")
            return False

    except Exception as e:
        logger.debug(f"Validation error: {e}")
        return False


def _detect_file_extension(file_data: bytes) -> str:
    """
    Detect file extension from file signature (magic bytes)

    Args:
        file_data: Binary file data

    Returns:
        File extension
    """
    if len(file_data) < 8:
        return '.bin'

    # Check for ZIP-based formats (PPTX, DOCX, etc.)
    if file_data[:4] == b'PK\x03\x04':
        # Further check for PPTX by looking for specific content
        if b'ppt/' in file_data[:1000] or b'[Content_Types].xml' in file_data[:2000]:
            return '.pptx'
        elif b'word/' in file_data[:1000]:
            return '.docx'
        elif b'xl/' in file_data[:1000]:
            return '.xlsx'
        return '.zip'

    # Check for old Office format (OLE2)
    if file_data[:8] == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1':
        # This is an OLE2 file, likely .ppt or .doc
        # Default to .ppt for PowerPoint context
        return '.ppt'

    return '.bin'


def extract_embedded_pptx_via_com(pptx_path: Path, output_dir: Path) -> List[Path]:
    """
    Extract embedded PowerPoint files using COM automation
    Works even for corrupted/non-standard PPTX files that PowerPoint can open

    Args:
        pptx_path: Path to the source PowerPoint file
        output_dir: Directory to save extracted files

    Returns:
        List of paths to extracted PowerPoint files
    """
    if not COM_AVAILABLE:
        logger.debug("COM not available, skipping COM extraction")
        return []

    pptx_path = Path(pptx_path).resolve()
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    extracted_files = []

    try:
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")

        try:
            # Hide window if possible
            powerpoint.Visible = False
        except:
            pass

        try:
            # Open presentation
            presentation = powerpoint.Presentations.Open(
                str(pptx_path),
                ReadOnly=True,
                Untitled=True,
                WithWindow=False
            )

            logger.info(f"COM: Opened presentation with {presentation.Slides.Count} slides")

            # Iterate through all shapes in all slides
            for slide_idx in range(1, presentation.Slides.Count + 1):
                slide = presentation.Slides(slide_idx)

                for shape_idx in range(1, slide.Shapes.Count + 1):
                    try:
                        shape = slide.Shapes(shape_idx)

                        # Check if shape has OLEFormat (embedded object)
                        if hasattr(shape, 'OLEFormat'):
                            try:
                                ole_format = shape.OLEFormat
                                prog_id = ole_format.ProgID

                                # Check if it's a PowerPoint object
                                if 'PowerPoint' in prog_id or 'Presentation' in prog_id:
                                    logger.info(f"COM: Found embedded PowerPoint in slide {slide_idx}, shape {shape_idx}")

                                    # Create temporary file to save the object
                                    temp_filename = f"embedded_s{slide_idx}_sh{shape_idx}.pptx"
                                    temp_path = output_dir / temp_filename

                                    # Try to activate and save the embedded object
                                    try:
                                        ole_format.DoVerb(1)  # Open for editing
                                        # Save the active presentation (the embedded one)
                                        embedded_pres = powerpoint.ActivePresentation
                                        embedded_pres.SaveAs(str(temp_path))
                                        embedded_pres.Close()

                                        extracted_files.append(temp_path)
                                        logger.info(f"COM: Extracted to {temp_filename}")

                                    except Exception as save_error:
                                        logger.debug(f"COM: Could not extract via DoVerb: {save_error}")

                            except Exception as ole_error:
                                # Shape doesn't have valid OLEFormat, skip
                                pass

                    except Exception as shape_error:
                        # Error accessing shape, skip
                        pass

            # Close presentation
            presentation.Close()

        finally:
            # Clean up PowerPoint
            powerpoint.Quit()
            pythoncom.CoUninitialize()

        logger.info(f"COM: Extracted {len(extracted_files)} embedded files")

    except Exception as e:
        logger.error(f"COM extraction failed: {e}")
        try:
            pythoncom.CoUninitialize()
        except:
            pass

    return extracted_files


def extract_embedded_pptx(pptx_path: Path, output_dir: Path) -> List[Path]:
    """
    Convenience function to extract all embedded PowerPoint files
    Uses multiple methods in order of reliability:
    1. ZIP extraction (fast, standard PPTX)
    2. COM extraction (works for corrupted/non-standard files)
    3. python-pptx OLE extraction (fallback)

    Args:
        pptx_path: Path to the source PowerPoint file
        output_dir: Directory to save extracted files

    Returns:
        List of paths to extracted PowerPoint files
    """
    pptx_path = Path(pptx_path)
    all_extracted = []

    # Method 1: Try ZIP extraction first (fastest for standard PPTX)
    logger.info(f"Extraction Method 1: ZIP-based extraction")
    try:
        zip_files = extract_embedded_pptx_from_zip(pptx_path, output_dir)
        if zip_files:
            logger.info(f"✓ ZIP method: Successfully extracted {len(zip_files)} files")
            all_extracted.extend(zip_files)
            return all_extracted
        else:
            logger.info("✗ ZIP method: No embedded files found")
    except zipfile.BadZipFile:
        logger.warning("✗ ZIP method: File is not a valid ZIP (corrupted or non-standard PPTX)")
    except Exception as e:
        logger.warning(f"✗ ZIP method failed: {e}")

    # Method 2: Try COM extraction (works for corrupted files)
    logger.info(f"Extraction Method 2: COM automation extraction")
    try:
        com_files = extract_embedded_pptx_via_com(pptx_path, output_dir)
        if com_files:
            logger.info(f"✓ COM method: Successfully extracted {len(com_files)} files")
            all_extracted.extend(com_files)
            return all_extracted
        else:
            logger.info("✗ COM method: No embedded files found")
    except Exception as e:
        logger.warning(f"✗ COM method failed: {e}")

    # Method 3: Fallback to python-pptx OLE extraction
    logger.info(f"Extraction Method 3: python-pptx OLE extraction")
    try:
        extractor = PPTXExtractor(pptx_path)
        extractor.load_presentation()
        ole_objects = extractor.extract_ole_objects()
        ole_files = extractor.save_ole_objects(ole_objects, output_dir)

        if ole_files:
            logger.info(f"✓ OLE method: Successfully extracted {len(ole_files)} files")
            all_extracted.extend(ole_files)
        else:
            logger.info("✗ OLE method: No embedded files found")
    except Exception as e:
        logger.warning(f"✗ OLE method failed: {e}")

    if not all_extracted:
        logger.warning(f"⚠ All extraction methods failed for {pptx_path.name}")

    return all_extracted
